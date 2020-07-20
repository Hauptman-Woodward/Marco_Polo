import pytest
import os
from PyQt5.QtTest import QTest
import sys
from PyQt5.QtWidgets import QApplication
from pathlib import Path
from random import randint

from polo import IMAGE_CLASSIFICATIONS
from polo.windows.pptx_dialog import PptxDesignerDialog
from polo.utils.io_utils import RunDeserializer, RunLinker
from pptx import Presentation


dirname = Path(os.path.dirname(__file__))
app = QApplication(sys.argv)


@pytest.fixture
def run():
    test_run = RunDeserializer(dirname.joinpath(
        'test_files/xtals/test_vis.xtal'))
    test_run = test_run.xtal_to_run()
    return test_run


@pytest.fixture
def write_bin():
    return dirname.joinpath('test_files/write_bin')


@pytest.fixture
def run_dict(run):
    return {run.run_name: run}


@pytest.fixture
def dummy_pptx_path(write_bin):
    return write_bin.joinpath('dummy.pptx')


@pytest.fixture
def run_random_xtals(run):
    # run with random images classified as crystals
    num_images = 20
    rand_indices = [randint(0, len(run.images)-1) for i in range(num_images)]
    for i in rand_indices:
        run.images[i].human_class = IMAGE_CLASSIFICATIONS[0]

    return run, num_images


@pytest.fixture
def all_xtal_files():
    xtal_dir = os.path.join(str(dirname), 'test_files/xtals')
    return [os.path.join(xtal_dir, x) for x in os.listdir(xtal_dir)]


@pytest.fixture
def all_runs_random_xtals(all_xtal_files):
    runs = []
    for run in all_xtal_files:
        deserial = RunDeserializer(run)
        runs.append(deserial.xtal_to_run())

    runs = RunLinker.the_big_link(runs)

    primary_run = runs.pop()
    num_images = 20
    rand_indices = [randint(0, len(primary_run.images)-1)
                    for i in range(num_images)]
    for i in rand_indices:
        primary_run.images[i].human_class = IMAGE_CLASSIFICATIONS[0]
    return primary_run, runs, num_images


@pytest.fixture
def pptx_human_xtals(run_dict, dummy_pptx_path):
    dummy_pptx_path = str(dummy_pptx_path)
    pptx = PptxDesignerDialog(runs=run_dict)
    pptx.ui.checkBox.setChecked(True)  # crystal classification
    pptx.ui.checkBox_6.setChecked(True)  # human classifier
    pptx.ui.lineEdit_3.setText(dummy_pptx_path)
    pptx.ui.lineEdit.setText(dummy_pptx_path)
    pptx.ui.lineEdit_2.setText('Test Write')

    return pptx


def test_image_checkboxes_dict(run_dict):
    pptx = PptxDesignerDialog(runs=run_dict)
    keys = list(pptx._image_class_checkboxes.keys())
    for k in keys:
        assert k in IMAGE_CLASSIFICATIONS


def test_pptx_designer_init(run_dict):
    pptx = PptxDesignerDialog(runs=run_dict)
    assert pptx
    assert isinstance(pptx, PptxDesignerDialog)


def test_write_pptx(pptx_human_xtals, run_random_xtals, dummy_pptx_path):
    # set to only write crystal classified marco images

    run_random_xtals, num_images = run_random_xtals
    dummy_pptx_path = str(dummy_pptx_path)
    if os.path.isfile(dummy_pptx_path):
        os.remove(dummy_pptx_path)

    assert not os.path.isfile(dummy_pptx_path)

    pptx_human_xtals._write_presentation(run=run_random_xtals)

    assert os.path.exists(dummy_pptx_path)
    assert os.path.isfile(dummy_pptx_path)

    pres = Presentation(dummy_pptx_path)
    assert pres.slides  # has slides
    assert len(pres.slides) == num_images + 1
    # should be one slide per image included plus the title slide


def test_write_pptx_all_specs(pptx_human_xtals, all_runs_random_xtals, 
                              dummy_pptx_path):
    primary_run, other_runs, num_images = all_runs_random_xtals
    dummy_pptx_path = str(dummy_pptx_path)
    if os.path.isfile(dummy_pptx_path):
        os.remove(dummy_pptx_path)

    assert not os.path.isfile(dummy_pptx_path)
    pptx_human_xtals.ui.checkBox_8.setChecked(True)
    pptx_human_xtals._write_presentation(primary_run)

    assert os.path.exists(dummy_pptx_path)
    assert os.path.isfile(dummy_pptx_path)

    pres = Presentation(dummy_pptx_path)

    num_crystal_images = len(primary_run.image_filter_query(
        ['Crystals'], human=True, marco=False, favorite=False))
    assert pres.slides
    assert len(pres.slides) != 1
    assert len(pres.slides) == (num_crystal_images * 2) + 1


def test_write_pptx_all_dates(pptx_human_xtals, all_runs_random_xtals,
                              dummy_pptx_path):
    primary_run, other_runs, num_images = all_runs_random_xtals
    dummy_pptx_path = str(dummy_pptx_path)
    if os.path.isfile(dummy_pptx_path):
        os.remove(dummy_pptx_path)

    assert not os.path.isfile(dummy_pptx_path)
    pptx_human_xtals.ui.checkBox_9.setChecked(True)
    pptx_human_xtals._write_presentation(primary_run)

    assert os.path.exists(dummy_pptx_path)
    assert os.path.isfile(dummy_pptx_path)

    pres = Presentation(dummy_pptx_path)

    num_crystal_images = len(primary_run.image_filter_query(
        ['Crystals'], human=True, marco=False, favorite=False))
    assert pres.slides
    assert len(pres.slides) != 1
    assert len(pres.slides) == (num_crystal_images * 2) + 1



    
