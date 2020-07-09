from pptx import Presentation
from pptx.util import Inches, Pt
import os
ims = '/home/ethan/Pictures/test_pics/deeper'

images = [os.path.join(ims, f) for f in os.listdir(ims)]

from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)



# add the title

title = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(10), Inches(1))
tf = title.text_frame
p = tf.add_paragraph()
p.text = 'Well 18 From X to X'
p.font.size = Pt(24)


# add the images here

image_size = round((10 - 2) / len(images), 1)
left, top = 1, 3
for image in images:
    slide.shapes.add_picture(image, Inches(left), Inches(top), height=Inches(image_size))
    label = slide.shapes.add_textbox(
        Inches(left), Inches(top + image_size) + Inches(0.5), Inches(image_size), Inches(1))
    label.rotation = 90
    tf = label.text_frame
    p = tf.add_paragraph()
    p.text = '12-31-1997'
    p.font.size = Pt(12)



    left += image_size

prs.save('test.pptx')
    