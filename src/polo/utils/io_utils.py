import base64
import csv
import inspect
import json
import os
import sys
import time
import xml.etree.ElementTree as ET
from datetime import datetime
from pathlib import Path

from jinja2 import Template
from PyQt5 import QtWidgets
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QBrush, QColor, QIcon, QImage, QPainter, QPixmap
from PyQt5.QtWidgets import QAction, QApplication, QGridLayout
from pptx import Presentation
from pptx.util import Inches, Pt

from polo import *
from polo.threads.thread import QuickThread
from polo.utils.dialog_utils import *
from polo.utils.exceptions import EmptyRunNameError, ForbiddenImageTypeError
from polo.utils.math_utils import *
from polo.utils.unrar_utils import *

logger = make_default_logger(__name__)


class RunImporter():

    allowed_filetypes = set(['.xtal', '.rar'])

    def __init__(self, target_path):
        self.target_path = Path(target_path)
        self.imported_run = None
        self.is_running = False
        logger.debug('Created {} for {}'.format(self, self.target_path))
    
    @staticmethod
    def make_xtal_file_dialog(parent=None):
        '''Create a file dialog specifically for browsing for
        xtal files.

        :param parent: Parent for the file dialog, defaults to None
        :type parent: QDialog, optional
        :return: QFileDialog
        :rtype: QFileDialog
        '''
        file_dlg = QtWidgets.QFileDialog(parent=parent)
        file_dlg.setFileMode(QtWidgets.QFileDialog.ExistingFiles)
        file_dlg.setNameFilter('xtal or xtals (*.xtal *.xtals)')
        return file_dlg
    
    @staticmethod
    def crack_open_a_rar_one(rar_path):
        '''Unrar a directory of images

        :param rar_path: Path to rar archive file
        :type rar_path: str or Path
        :return: Unrar result
        :rtype: str, exception or int
        '''

        parent_path = rar_path.parent
        return unrar_archive(rar_path, parent_path)
    
    def _target_is_valid(self):
        '''Private method to check if a the file or directory referenced by
        :attr:`target_path` is valid for import.

        :return: True if could be imported, False otherwise
        :rtype: bool
        '''
        if self.target_path.exists():
            if (self.target_path.is_file()
                and self.target_path.suffix in RunImporter.allowed_filetypes):
                return True
            elif (self.target_path.is_dir() 
                 and list_dir_abs(str(self.target_path), allowed=True)):
                return True
        return False
    
    def create_import_thread(self):
        '''Creates a :class:`polo.threads.thread.QuickThread` with function
        based on the type of file that is to be imported.

        :return: QuickThread
        :rtype: QuickThread
        '''
        if self._target_is_valid():
            # check if needs unrar
            
            if self.target_path.suffix == '.rar':
                import_thread = QuickThread(
                    RunImporter.crack_open_a_rar_one, rar_path=self.target_path)
            elif self.target_path.suffix == '.xtal':
                self._xtal_reader = RunDeserializer(str(self.target_path))
                import_thread = QuickThread(
                    self._xtal_reader.xtal_to_run
                )
            else:
                # is a directory create a dummy thread
                import_thread = QuickThread(
                    lambda: self.target_path
                )
            return import_thread


class RunSerializer():

    def __init__(self, run=None):
        self.run = run

    @classmethod
    def make_thread(cls, job_function, **kwargs):
        '''
        Creates a new :class:`QuickThread` object. The job function is the
        function the thread will execute and and arguments that the job
        function requires should be passed has keyword arguments. These are
        stored as a dictionary in the new thread object until the thread is
        activated and they are passed as arguments.
        '''
        return QuickThread(job_function, parent=None, **kwargs)

    @staticmethod
    def path_suffix_checker(path, desired_suffix):
        '''
        Check is a file path has a desired suffix, if not then replace the
        current suffix with the desired suffix. Useful for checking filenames
        that are taken from user input.

        :param desired_suffix: File extension for given file path.
        '''
        try:
            if desired_suffix:
                if isinstance(path, str):
                    path = Path(path)
                if path.suffix == desired_suffix:
                    return str(path)
                else:
                    return str(path.with_suffix(desired_suffix))
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                e, self.path_suffix_checker))
            return None

    @staticmethod
    def path_validator(path, parent=False):
        '''
        Tests to ensure a path exists. Passing parent = True will check for
        the existence of the parent directory of the path.
        '''
        if isinstance(path, (str, Path)):
            if isinstance(path, str):
                path = Path(path)
            if parent:
                return path.parent.exists()
            else:
                return path.exists()
        else:
            return False

    def __repr__(self):
        s = ''
        for key, value in self.__dict__.items():
            s += '{}: {}.unitsn'.format(key, value)
        return s[:-1]


class HtmlWriter(RunSerializer):

    def __init__(self, run, **kwargs):
        super(HtmlWriter, self).__init__(run)
        self.__dict__.update(kwargs)

    @staticmethod
    def make_template(template_path):
        '''Given a path to an html file to serve as a jinja2 template, read the
        file and create a new template object.

        :param template_path: Path to the jinja2 template file.
        :type temlate_path: str
        '''
        if isinstance(template_path, Path):
            template_path = str(template_path)

        with open(template_path, 'r') as template:
            contents = template.read()
            logger.debug('Read Jinja template at {}'.format(template_path))
            return Template(contents)

    def write_complete_run(self, output_path, encode_images=True):
        '''Create an HTML report from a :class:`Run` or :class:`HWIRun`
        instance.

        :param output_path: Path to write html file to.
        :type output_path: str or Path
        :param encode_images: Write images as base64 directly to the html file,
                              defaults to True. Greatly increases the file size
                              but means that report will still contain images
                              even if the originals are deleted or removed.
        :type encode_images: bool, optional
        :return: Path to html report if write succeeds, Exception otherwise. 
        :rtype: str or Exception
        '''
        try:
            if HtmlWriter.path_validator(output_path, parent=True) and self.run:
                output_path = HtmlWriter.path_suffix_checker(output_path, '.html')
                if output_path:
                    if encode_images:
                        self.run.encode_images_to_base64()
                    run = XtalWriter.clean_run_for_save(self.run)
                    images = json.loads(json.dumps(
                        run.images, default=XtalWriter.json_encoder))
                    [RunDeserializer.clean_base64_string(image['_bites'], str) for image in images]

                    template = HtmlWriter.make_template(RUN_HTML_TEMPLATE)
                    if template:
                        html = template.render(
                            images=images, run_name=self.run.run_name,
                            annotations='No annotations')
                        with open(output_path, 'w') as html_file:
                            html_file.write(html)
                            logger.debug('Wrote HTML report to {}'.format(output_path))
                            return True
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                e, self.write_complete_run))
            return e

    def write_grid_screen(self, output_path, plate_list, well_number,
                          x_reagent, y_reagent, well_volume, run_name=None):
        '''Write the contents of optimization grid screen to an html file.

        :param output_path: Path to html file
        :type output_path: str
        :param plate_list: list containing grid screen data
        :type plate_list: list
        :param well_number: well number of hit screen is created from
        :type well_number: int or str
        :param x_reagent: reagent varied in x direction
        :type x_reagent: Reagent
        :param y_reagent: Reagent varied in y direction
        :type y_reagent: Reagent
        :param well_volume: Volume of well used in screen
        :type well_volume: int or str
        :param run_name: name of run, defaults to None
        :type run_name: str, optional
        '''
        try:
            if HtmlWriter.path_validator(output_path, parent=True):
                output_path = HtmlWriter.path_suffix_checker(output_path, '.html')
                template = HtmlWriter.make_template(SCREEN_HTML_TEMPLATE)
                if not run_name:
                    run_name = self.run.run_name
                html = template.render(plate_list=plate_list, run_name=run_name,
                                    well_number=well_number, date=datetime.now(),
                                    x_reagent_stock=x_reagent,
                                    y_reagent_stock=y_reagent,
                                    well_volume=str(well_volume))
                with open(output_path, 'w') as screen_html:
                    screen_html.write(html)
                    logger.debug('Wrote grid screen report to {}'.format(output_path))
        except Exception as e:
            logger.error('Caught {} calling {}'.format(e, self.write_grid_screen))
            make_message_box(
                parent=None,
                message='Failed to write screen to {}. {}'.format(
                    output_path, e
                )
            ).exec_()


class RunCsvWriter(RunSerializer):

    def __init__(self, run, output_path=None, **kwargs):
        self.__dict__.update(kwargs)
        self.output_path = output_path
        super(RunCsvWriter, self).__init__(run)

    @classmethod
    def image_to_row(cls, image):
        '''Given an :class:`Image` object, convert it into a list that could be
        easily written to a csv file.

        :param image: :class:`Image` object to convert to list
        :type image: Image
        :return: List
        :rtype: list
        '''
        row = {}
        for attr, value in image.__dict__.items():
            if attr[0] == '_':  # remove _ from hidden attributes so looks nice when displayed
                attr = attr[1:]
            if isinstance(value, Cocktail):
                row[attr] = value.number
            elif isinstance(value, Image):
                row[attr] = value.path
            elif isinstance(value, dict):
                # unwrap the dict
                for attr_b, value_b in value.items():
                    if attr_b not in row:  # only add if will not override higher level attr
                        # only go one level deep for now
                        row[attr_b] = str(value_b)
            elif isinstance(value, bytes) or attr == 'bites':
                continue  # currently do not encode bytes (base 64 stuff)  
            else:
                row[attr] = str(value)  # default to case to string
        return row

    @property
    def output_path(self):
        '''Get the hidden attribute `_output_path`.

        :return: Output path
        :rtype: str
        '''
        return self._output_path

    @property
    def fieldnames(self):  # could use more efficent way of determining feildnames
        '''Get the current fieldnames based on the data stored in the
        :attr:`~polo.utils.io_utils.RunCsvWriter.run` attribute.
        Currently is somewhat expensive to call since it
        requires parsing all records in 
        :attr:`~polo.utils.io_utils.RunCsvWriter.run`
        in order to determine all the
        fieldnames that should be included in order to definitely avoid
        keyerrors later down the line.

        :return: List of fieldnames (headers) for the csv data 
        :rtype: list
        '''
        rows = [row for row in self]
        fieldnames = set([])
        for row in rows:
            fieldnames = fieldnames.union(set(row.keys()))

        return fieldnames

    @output_path.setter
    def output_path(self, new_path):
        if isinstance(new_path, Path):
            new_path = str(new_path)

        self._output_path = new_path

    def get_csv_data(self):
        '''Convert the  :attr:`~polo.utils.io_utils.RunCsvWriter.run`
        attribute to csv style data. Returns a tuple of
        headers and a list of dictionaries with each dictionary representing
        one row of csv data.

        :return: Tuple, list of headers and list of dicts
        :rtype: tuple
        '''
        try:
            rows = [row for row in self]
            fieldnames = set([])
            for row in rows:
                fieldnames = fieldnames.union(set(row.keys()))
            return fieldnames, rows
        except Exception as e:
            logger.error('Caught and re-raised {} while calling {}'.format(
                            e, self.get_csv_data))
            raise e

    def write_csv(self):
        '''Write the :class:`Run` object referenced by the 
        :attr:`~polo.utils.io_utils.RunCsvWriter.run` 
        attribute as a csv file to the location specified
        by the  :attr:`~polo.utils.io_utils.RunCsvWriter.output_path`
        attribute. 

        :return: True, if csv file content was written successfully,
                 return error thrown otherwise.
        :rtype: Bool or Exception
        '''
        try:
            rows = [row for row in self]
            fieldnames = set([])
            for row in rows:
                fieldnames = fieldnames.union(set(row.keys()))
            with open(self.output_path, 'w', newline='') as csv_path:
                writer = csv.DictWriter(csv_path, fieldnames)
                writer.writeheader()
                for row in rows:
                    writer.writerow(row)
            logger.debug('Wrote csv data to {}'.format(self.output_path))
            return True
        except Exception as e:
            logger.warning(
                'Caught exception {} at {}'.format(e, self.write_csv))
            return e

    def __iter__(self):
        for image in self.run.images:
            yield RunCsvWriter.image_to_row(image)

class SceneExporter():

    def __init__(self, graphics_scene=None, file_path=None):
        self.graphics_scene = graphics_scene
        self.file_path = file_path
    
    @staticmethod
    def write_image(scene, file_path):
        '''Write the contents of a :class:`QGraphicsScene`
        to a png file.

        :param scene: :class:`QGraphicsScene` to convert to image file.
        :type scene: QGraphicsScene
        :param file_path: Path to save image to.
        :type file_path: str
        :return: File path to saved image if successful, Exception otherwise.
        :rtype: str or Exception
        '''
        try:
            image = QImage(
                scene.width(), scene.height(),
                QImage.Format_ARGB32_Premultiplied
                )
            painter = QPainter(image)
            scene.render(painter)
            image.save(file_path)
            painter.end()
            return file_path
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                            e, SceneExporter.write_image))
            return e
    
    def write(self):
        return SceneExporter.write_image(self.graphics_scene, self.file_path)


class MsoWriter(RunSerializer):

    mso_version = 'msoversion2'

    def __init__(self, run, output_path):
        self.run = run
        self.output_path = output_path

    @staticmethod
    def row_formater(cocktail_row):
        '''Format a cocktail row as read from a cocktail csv
        file to an mso file row. Main change is appending empty
        strings to the cocktail row so list ends up always having
        a length of 17. This is important because the image
        classification code always occurs at the 18th item in
        an mso file row.

        :param cocktail_row: Cocktail row as read from cocktail csv file.
        :type cocktail_row: list
        :return: Cocktail row reformated for mso writing.
        :rtype: list
        '''
        # well number should be first index in the list
        # total list length should be 18 last item is the mso code
        return cocktail_row + ([''] * (len(cocktail_row) - 17))

    
    @property
    def first_line(self):
        '''Create the first line of the mso file.

        :return: List to write as the first line of the mso file.
        :rtype: list
        '''
        return [
            MsoWriter.path_suffix_checker(self.run.run_name, '.rar'), 
            self.mso_version
        ]
    
    def get_cocktail_csv_data(self):
        '''Reads and returns the cocktail csv data assigned
        to the  :attr:`~polo.utils.io_utils.MsoWriter.cocktail_menu`
        attribute of the MsoWriter's
        :attr:`~polo.utils.io_utils.MsoWriter.run` attribute.

        :return: List of lists containing cocktail csv data.
        :rtype: list
        '''
        cocktail_menu = self.run.cocktail_menu.path
        with open(cocktail_menu, 'r') as menu_file:
            next(menu_file)  # skip first row
            return [row for row in csv.reader(menu_file)]

    def write_mso_file(self, use_marco_classifications=False):
        '''Writes an mso formated file for use with MarcoScopeJ based on
        the images and classifications of the :class:`Run` instance
        referenced by the MsoWriter's :attr:`~polo.utils.io_utils.MsoWriter.run`
        attribute.

        :param use_marco_classifications: Include the MARCO classification
                                          in the mso file instead of human
                                          classifications, defaults to False
        :type use_marco_classifications: bool, optional
        :return: True if file is written successfully, False otherwise.
        :rtype: bool
        '''
        if isinstance(self.run, HWIRun):  # must be hwi run to write to mso
            cocktail_data = self.get_cocktail_csv_data()
            self.output_path = str(MsoWriter.path_suffix_checker(self.output_path, '.mso'))
            with open(self.output_path, 'w',  newline='') as mso_file:
                # newline='' added for windows compatibility
                # see https://stackoverflow.com/questions/3348460/csv-file-written-with-
                # python-has-blank-lines-between-each-row
                writer = csv.writer(mso_file, delimiter='\t')
                first_line = [str(i).strip() for i in self.first_line]
                header = [str(i).strip() for i in cocktail_data.pop(0)]
                writer.writerow(first_line)
                writer.writerow(header)
                for row in cocktail_data:
                    row = MsoWriter.row_formater(row)
                    well_num = int(float(row[0]))

                    image = self.run.images[well_num-1]
                    if image and image.human_class:
                        row.append(MSO_DICT[image.human_class])
                        writer.writerow(row)
                    elif use_marco_classifications and image.machine_class:
                        row.append(MSO_DICT[image.machine_class])
                        writer.writerow(row)
            return True
        else:
            return False


class MsoReader():
    '''The MsoReader class is used to parse the content of mso formated
    files and apply the image classifications stored in these files to
    runs in Polo.
    '''

    def __init__(self, mso_path):
        self.mso_path = mso_path
    
    
    @staticmethod
    def read_mso_classification(mso_classification):
        '''Helper method to read and convert image classifications in a mso
        file to MARCO image classifications. The exact conversion scheme is
        layed out in the :const:`REV_MSO_DICT` constant.
        Additionally, MarcoscopeJ will allow
        images to have multiple classifications by assigning multiple
        image codes seperated by "-". If this is the case Polo takes the
        classification corresponding to the mso code with the highest value.

        :param mso_classification: Mso image classification code
        :type mso_classification: str
        :return: MARCO classification if code can be decoded, None otherwise
        :rtype: str or None
        '''
        try:
            mso_classification = mso_classification.split('-')

            mso_codes = [int(code.strip()) for code in mso_classification 
                        if int(code.strip()) in REV_MSO_DICT]
            if mso_codes:
                return REV_MSO_DICT[max(mso_codes)]
            else:
                return None
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                            e, MsoReader.read_mso_classification))
            return None  # assume no mso classification if throws an error

    def classify_images_from_mso_file(self, images):
        '''Applies the image classifications stored in an mso file to a
        collection of :class:`~polo.crystallography.image.Image` objects.
        Allows for some degree of compatability
        with MarcoscopeJ and for users who have stored their image classifications
        in the mso format. Additionally, human classification backups are
        saved as mso files when Polo is closed as they take up much less
        space than xtal files.

        :param images: List of images to apply classifications to
        :type images: list
        :return: List of images with mso classifications applied,
                 or Exception if this fails
        :rtype: list or Exception
        ''' 
        try:
            with open(str(self.mso_path)) as mso:
                reader = csv.reader(mso, delimiter='\t')
                next(reader)
                next(reader)  # skip first two header lines
                if reader:
                    for row in reader:
                        well_index = int(row[0]) - 1
                        classification = MsoReader.read_mso_classification(row[-1])
                        try:
                            images[well_index].human_class = classification
                        except IndexError:
                            continue
                return images
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                            e, self.classify_images_from_mso_file))
            return e

class JsonWriter(RunSerializer):
    '''Small class that can be used to serialize a run to a
    json formated file.

    :param run: Run to write as a json file
    :type run: Run or HWIRun
    :param output_path: Path to write json file to
    :type output_path: str or Path
    '''
    def __init__(self, run, output_path):
        self.run = run
        self.output_path = output_path
    
    @staticmethod
    def json_encoder(obj):
        '''General purpose json encoder for encoding python objects. Very
        similar to the encoder function 
        :meth:`~polo.utils.io_utils.XtalWriter.json_encoder` except does not
        include class and module information in the returned dictionary. If
        the object cannot be converted to a dictionary it is returned as a
        string.   

        :param obj: Object to convert to dictionary
        :type obj: obj
        :return: dict or str
        :rtype: dict or str
        '''
        d = None
        if hasattr(obj, '__dict__'):  # can send to dict object
            d = obj.__dict__
        else:  # not castable to dict
            if isinstance(obj, bytes):  # likely the base64 encoded image
                d = obj.decode('utf-8')
            else:
                d = str(obj)  # if all else fails case to string
        return d
    
    def write_json(self):
        '''Write the :class:`Run` instance referenced by the :attr:`~polo.utils.io_utils.JsonWriter.run`
        attribute to a json file at
        the location specified by the
        :attr:`~polo.utils.io_utils.JsonWriter.output_path` attribute.
        If the file is written successfully returns True
        otherwise returns an Exception.

        :return: True or Exception
        :rtype: bool, Exception
        '''
        try:
            clean_run = XtalWriter.clean_run_for_save(self.run)
            json_content = json.dumps(clean_run, ensure_ascii=True, indent=4,
                                  default=XtalWriter.json_encoder)
            output_path = JsonWriter.path_suffix_checker(
                str(self.output_path), '.json')
            with open(output_path, 'w') as json_file:
                json_file.write(json_content)
            return True
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                            e, self.write_json))
            return e


class XtalWriter(RunSerializer):
    header_flag = '<>'  # do not change unless very good reason
    header_line = '{}{}:{}\n'
    file_ext = '.xtal'

    def __init__(self, run, main_window, **kwargs):
        self.__dict__.update(kwargs)
        self.main_window = main_window
        super(XtalWriter, self).__init__(run)

    @property
    def xtal_header(self):
        '''Creates the header for an xtal file when called. Header lines are
        indicated as such by the string in the header_line constant,
        which should be '<>'. The last line of the header will be a row
        of equal signs and then the actual json content begins on the
        next line.
        '''
        header = ''
        header += self.header_line.format(
            self.header_flag, 'SAVE TIME', datetime.now())
        header += self.header_line.format(
            self.header_flag, 'VERSION', polo_version)
        for key, value in self.__dict__.items():
            header += self.header_line.format(
                self.header_flag, str(key).upper(), value)
        return header + '='*79 + '\n'
        # add ==== as a break between json data and header data

    @staticmethod
    def json_encoder(obj):
        '''Use instead of the default json encoder when writing an xtal file. 
        If the encoded object is from a module within Polo will include a 
        module and class identifier so it can be more easily deserialized 
        when loaded back into the program.

        :param: obj: An object to serialize to json.

        :returns: A dictionary or string version of the passed object
        '''
        d = None
        if hasattr(obj, '__dict__'):  # can send to dict object
            d = obj.__dict__
            d['__class__'] = obj.__class__.__name__
            d['__module__'] = obj.__module__
            # store module and class name along with object as dict
        else:  # not castable to dict
            if isinstance(obj, bytes):  # likely the base64 encoded image
                d = obj.decode('utf-8')
            else:
                d = str(obj)  # if all else fails case to string
        return d

    @staticmethod
    def clean_run_for_save(run):
        '''Remove circular references from the run passed through the `run`
        argument to avoid issues when writing to json files.

        :param run: Run to clean (remove circular references)
        :type run: Run or HWIRun
        :return: Run, free from circular references
        :rtype: Run or HWIRun
        '''

        if run:
            run.previous_run, run.next_run, run.alt_spectrum = (
                None, None, None)
            for image in run.images:
                if image:
                    image.previous_image, image.next_image, image.alt_image = (
                        None, None, None
                    )
        return run

    def write_xtal_file_on_thread(self, output_path):
        """Wrapper method around :meth:`~polo.utils.io_utils.XtalWriter.write_xtal_file` 
        that executes on a :class:`QuickThread` instance to prevent freezing the 
        GUI when saving large xtal files

        :param output_path: Path to xtal file
        :type output_path: str
        """

        self.thread = XtalWriter.make_thread(
            self.write_xtal_file, output_path=output_path)
        # connect to something when thread finishes
        self.thread.finished.connect(self.finished_save)
        QApplication.setOverrideCursor(Qt.WaitCursor)
        self.main_window.setEnabled(False)
        self.thread.start()

    def finished_save(self):
        self.main_window.setEnabled(True)
        QApplication.restoreOverrideCursor()  # return cursor to normal
        # should contain path to now saved file
        if os.path.exists(str(self.thread.result)):
            message = 'Saved current run to {}!'.format(self.thread.result)
        else:
            message = 'Failed to save run.{}'.format(self.thread.result)

        make_message_box(message=message).exec_()

    def write_xtal_file(self, output_path=None):
        '''Method to serialize run object to xtal file format.

        :param output_path: Xtal file path
        :type output_path: str
        :return: path to xtal file
        :rtype: str
        '''
        try:
            if XtalWriter.path_validator(output_path, parent=True):
                # path is good to go and ready to write file into
                self.run.encode_images_to_base64()
                run_str = self.run_to_dict()
                if isinstance(run_str, str):  # encoding worked, no errors caught
                    output_path = XtalWriter.path_suffix_checker(
                        output_path, self.file_ext)  # make sure has .xtal suffix
                
                    with open(str(output_path), 'w') as xtal_file:
                        xtal_file.write(self.xtal_header)
                        xtal_file.write(run_str)
                        return output_path
        except Exception as e:
            logger.error('Caught {} at {}'.format(
                e, self.write_xtal_file))
            return e

    def run_to_dict(self):
        '''Create a json string from the run stored in the run attribute.

        :return: Run instance serialized to json
        :rtype: str
        '''
        if self.run:
            try:
                clean_run = XtalWriter.clean_run_for_save(self.run)
                clean_run.encode_images_to_base64()
                return json.dumps(clean_run, ensure_ascii=True, indent=4,
                                  default=XtalWriter.json_encoder)
            except Exception as e:
                logger.error('Failed to encode {} to dict. Gave {}'.format(
                    self.run, e))
                return e

class XtalToImages():
    '''One of the features of Polo is the ability to save images, classifications
    and cocktail data all as one (xtal) file. This also means that users may
    be loading in their run data from xtal files but want some way to recover
    the images that are encoded (as base64) within the delicious onion like
    layers of the JSON formated xtal file. This class acomplishes this goal by
    intializing with a run that has base64 encoded image data and writes that
    data as jpegs to a user specified directory.

    Currently under construction as of 12-14-2020
    '''

    def __init__(self, run, output_dir):
        self.run = run
        self.output_dir = output_dir
    
    def write_images(self):
        for i in self._collect_writable_image_indicies():
            image_write_path = os.path.join(output_dir, os.path.basename(image.path))
    

    def _collect_writable_image_indicies(self):
        return [i for i in range(len(self.run.images)) if image.bites]
        



class RunDeserializer():  # convert saved file into a run

    def __init__(self, xtal_path):
        self.xtal_path = xtal_path

    @staticmethod
    def clean_base64_string(string, out_fmt=bytes):
        '''Image instances may contain byte strings that store their actual
        crystallization image encoded as base64. Previously, these byte strings
        were written directly into the json file as strings causing the b'
        byte string identifier to be written along with the actual base64 data.
        This method removes those artifacts if they are present and returns a
        clean byte string with only the actual base64 data.

        :param string: a string to interrogate
        :type string: str
        :return: byte string with non-data artifacts removed
        :rtype: bytes
        '''
        if string:
            if isinstance(string, bytes):
                string = str(string, 'utf-8')
            if string[0] == 'b':  # bytes string written directly to string
                string = string[1:]
            if string[-1] == "'":
                string = string[:-1]
            if string: 
                if isinstance(out_fmt, bytes):
                    return bytes(string, 'utf-8')
                else:
                    return string
            
    @staticmethod
    def dict_to_obj(d):
        '''Opposite of the obj_to_dict method in XtalWriter class, this method
        takes a dictionary instance that has been previously serialized and
        attempts to convert it back into an object instance. Used as the
        `object_hook` argument when calling `json.loads` to read xtal files.

        :param d: dictionary to convert back to object
        :type d: dict
        :return: an object
        :rtype: object
        '''
        if d:
            if '__class__' in d:  # is a serialized object
                class_name, mod_name = d.pop('__class__'), d.pop('__module__')
                try:
                    module = __import__(mod_name)
                    # BUG module is polo from init so if currently if classes are
                    # not imported into init file they will not be found here
                    class_ = getattr(module, class_name)
                except AttributeError as e:
                    logger.error('Caught {} while calling {}'.format(
                            e, RunDeserializer.dict_to_obj))
                    return None
                temp_d = {}
                for key, item in d.items():
                    if '__' in key:  # deal with private object properties
                        key = key.split('__')[-1]
                    elif '_' == key[0]:
                        key = key[1:]
                    temp_d[key] = item
                obj = class_(**temp_d)

                if isinstance(obj, Image):  # clean up base64 encoded data
                    if obj.bites:
                        obj.bites = RunDeserializer.clean_base64_string(
                            obj.bites)
            else:
                obj = d  # just a regular dictionary to read in
            return obj
        else:
            return None

    def xtal_header_reader(self, xtal_file_io):
        '''Reads the header section of an open xtal file. Should always be
        called before reading the json content of an xtal file. Note than
        xtal files must always have a line of equal signs before the json
        content even if there is no header content otherwise this method will
        read one line into the json content causing the json reader to
        throw an error.

        :param xtal_file_io: xtal file currently being read
        :type xtal_file_io: TextIoWrapper
        :return: xtal header contents
        :rtype: list
        '''
        # pulls info in the header of an xtal file
        header_data = [xtal_file_io.readline()]
        while header_data[-1][0:len(XtalWriter.header_flag)] == XtalWriter.header_flag:
            header_data.append(xtal_file_io.readline())
        return header_data

    def xtal_to_run(self, **kwargs):
        '''Attempt to convert the file specified by the path stored in the
        :attr:`~polo.utils.io_utils.RunDeserializer.xtal_path`
        attribute to a :class:`Run` object. 

        :return: Run object encoded by an xtal file
        :rtype: Run
        '''
        try:
            if 'xtal_path' in kwargs:
                xtal_path = kwargs['xtal_path']
            else:
                xtal_path = self.xtal_path
            xtal_path = str(xtal_path)
            if os.path.isfile(xtal_path):
                with open(xtal_path) as xtal_data:
                    header_data = self.xtal_header_reader(
                        xtal_data)  # must read header first
                    r = json.load(xtal_data,  # update date since datetime goes right to string
                                object_hook=RunDeserializer.dict_to_obj)
                    r.date = BarTender.datetime_converter(r.date)
                    r.save_file_path = xtal_path
                    return r
            else:
                return FileNotFoundError
        except Exception as e:
            logger.error('Caught {} while calling {}'.format(
                            e, self.xtal_to_run))
            return e


class PptxWriter():
    '''Use for creating pptx presentation slides from :class:`Run`
    or :class:`HWIRun` instances.

    :param output_path: Path to write pptx file to.
    :type output_path: str or Path
    :param included_attributes: [description], defaults to {}
    :type included_attributes: dict, optional
    :param image_types: Images included in the presentation
                        must have a classification in this set, defaults to None
    :type image_types: set or list, optional
    :param human: Use human classification as the image classification, defaults to False
    :type human: bool, optional
    :param marco: Use the MARCO classification as the image classification, defaults to False
    :type marco: bool, optional
    :param favorite: Only include images marked as favorite, defaults to False
    :type favorite: bool, optional
    '''
    slide_title_formater = 'Well Number {}'

    # 13.33 x 7.5 
    def __init__(self, output_path):
        self.output_path = output_path
        self._temp_images = []
        self._bumper = 1
        self._slide_width = 10
        self._slide_height = 6
        self._presentation = Presentation()
    
    def make_presentation(self, run, images, title, subtitle=None,
                                    cocktail_data=True, all_specs=False,
                                    all_dates=False):
        '''Create a pptx presentation file from a screening run instance. This
        should be the only public method in the class. All other methods are
        internal to actually creating the presentation file. 

        :param run: Run to create the presentation from
        :type run: Run
        :param title: Title for the presentation
        :type title: str
        :param subtitle: Subtitle for the presentation, defaults to None
        :type subtitle: str, optional
        :param cocktail_data: If True include cocktail data on slide, defaults to True
        :type cocktail_data: bool, optional
        :param all_specs: If True include all spectrum image slide, defaults to False
        :type all_specs: bool, optional
        :param all_dates: If True include all dates slide, defaults to False
        :type all_dates: bool, optional
        :return: True if presentation was created successfully, Exception otherwise
        :rtype: bool or Exception
        '''

        try:
            title_slide = self._add_new_slide(0)
            title_slide.shapes.title.text = title
            if subtitle:
                title_slide.placeholders[1].text = subtitle
            
            metadata = self._metadata_from_images(images, cocktail_data)
            self._add_images_and_metadata_to_slideshow(images, metadata, all_specs, all_dates)            
            self.output_path = RunSerializer.path_suffix_checker(
                                                                self.output_path, 
                                                                '.pptx'
                                                                )
            if images:  # make sure actualy were images to write to file
                self._presentation.save(str(self.output_path))
                return True
            else:
                return False
        except Exception as e:
            raise e
            logger.error('Caught {} while calling {}'.format(
                            e, self.make_presentation))
            return e
    
    def _add_images_and_metadata_to_slideshow(self, images, metadata, all_specs, all_dates):
        for image, image_metadata in zip(images, metadata):
            self._add_single_image_slide(
                image, PptxWriter.slide_title_formater.format(image.well_number),
                metadata=image_metadata
            )
            if all_specs and image.alt_image:
                self._add_multi_spectrum_slide(
                    image.get_linked_images_by_spectrum(), image.well_number
                )
            if all_dates and (image.next_image or image.previous_image):
                self._add_timeline_slide(image.get_linked_images_by_date(), 
                                        image.well_number
                                        )
            
    def _metadata_from_images(self, images, cocktail_data):
        all_metadata = []
        for image in images:
            image_metadata = str(image)
            if cocktail_data and hasattr(image, 'cocktail'):
                # user wants to display cocktail data and the image has it
                image_metadata += '\n\n' + str(image.cocktail)
            all_metadata.append(image_metadata)
        return all_metadata
    
    def _select_images(self, run):
        '''Selects images to add to the presentation.
        '''
        if self.images_indices:
            # Use these index instead of classifications
            return [run.images[i] for i in self.images_indices]
        else:
            return [image for image in run.images 
                    if image.standard_filter(
                        self.image_types, self.human, self.marco, self.favorite
                        )
                    ]
    
    def _delete_presentation(self):
        self._presentation = Presentation()
    
    def _delete_temp_images(self):
        '''Delete an temporary images used to create the pptx presentation.

        :return: True, if images are removed successfully, Exception otherwise.
        :rtype: bool or Exception
        '''
        try:
            [os.remove(img_path) for img_path in self._temp_images]
            return True
        except (FileNotFoundError, IsADirectoryError, PermissionError) as e:
            logger.error('Caught {} while calling {}'.format(
                            e, self._delete_temp_images))
            return e
    
    def sort_runs_by_spectrum(self, runs):
        '''Divids runs into two lists, one containing visible spectrum
        runs and another containing all non-visible runs.

        :param runs: List or runs
        :type runs: list
        :return: tuple, first item is visible runs second is non-visible runs
        :rtype: tuple
        '''
        visible, other = [], []
        for run in runs:
            if run.image_spectrum == IMAGE_SPECS[0]:
                visible.append(run)
            else:
                other.append(run)
        return visible, other
                
    def _add_new_slide(self, template=5):
        '''Add a new slide to the presentation referenced by the
        :attr:`_presentation` attribute.

        :param template: Slide template integer identifier, defaults to 5.
                         See the python-pptx package for more details on 
                         template integers
        :type template: int, optional
        :return: Presentation with the new slide added 
        :rtype: Presentation
        '''
        return self._presentation.slides.add_slide(
            self._presentation.slide_layouts[template])
    
    def _add_timeline_slide(self, images, well_number):
        '''Create a timeline (time resolved) slide that
        show the progression of a sample in a particular
        well.

        :param images: List of images to include in the slide
        :type images: list
        :param well_number: Well number to use in the title of the slide
        :type well_number: int
        :return: New slide
        :rtype: slide
        '''
        date_images = sorted(images, key=lambda i: i.date)
        new_slide = self._add_new_slide()
        labeler = lambda i: i.formated_date
        self._add_multi_image_slide(new_slide, images, labeler)
        title = 'Well {}: {} - {}'.format(
            well_number, date_images[0].formated_date, date_images[-1].formated_date)
        new_slide.shapes.title.text = title

        return new_slide

    def _add_multi_spectrum_slide(self, images, well_number):
        '''Create a slide to show a all spectrums a well has been
        imaged in.

        :param images: Images to include on the slide
        :type images: list
        :param well_number: Well number to use in the slide title
        :type well_number: int
        :return: New slide
        :rtype: slide
        '''
        spec_images = sorted(images, key=lambda i: len(i.spectrum))
        new_slide = self._add_new_slide()
        labeler = lambda i: i.spectrum
        self._add_multi_image_slide(new_slide, images, labeler)
        title = 'Well {} Alternate Imagers'.format(well_number)
        new_slide.shapes.title.text = title

        return new_slide
    
    def add_cocktail_slide(self, well, cocktail):
        '''Add slide with details on :class:`Cocktail` information.

        :param well: Well number to use in slide title
        :type well: int
        :param cocktail: Cocktail to write as a slide
        :type cocktail: Cocktail
        '''
        new_slide = self._add_new_slide(5)
        title = 'Well {} Cocktail: {}'.format(well, cocktail.number)
        new_slide.shapes.title.text = title

    def add_table_to_slide(self, slide, data, left, top):
        '''General helper method for adding a table to a slide.

        :param slide: Slide to add the table to
        :type slide: slide
        :param data: List of lists that has the data to write to the table
        :type data: list
        :param left: Left offset in inches to place to table
        :type left: float
        :param top: Top cordinate for placing the table
        :type top: float
        :return: Slide with table added
        :rtype: slide
        '''
        rows, cols = len(data), max([len(r) for r in data])
        shapes = slide.shapes
        
        width = (self.slide_width - (self._bumper * 2))
        height = self._slide_height - 2
        table = shapes.add_table(rows, cols, Inches(left), Inches(top), 
                                 Inches(width), Inches(height))

        for k in range(cols):
            table.columns[k].width = (self.slide_size - (self._bumper * 2)) / cols
            # set column width
        for i in range(rows):
            for j in range(cols):
                table.cell(i, j).text = data[i][j]
        return slide

    def _add_multi_image_slide(self, slide, images, labeler):
        '''Private method for adding a slide that will have multiple
        images.

        :param slide: Slide to add the images to 
        :type slide: slide
        :param images: Images to add to the slide
        :type images: list
        :param labeler: Function to use to label the individual images
        :type labeler: func
        :return: slide with images added
        :rtype: slide
        '''
        top = 2.5
        img_size = (self._slide_width - (self._bumper * 2)) / len(images)
        if img_size >= 0.4 * self._slide_height: img_size = 0.4 * self._slide_height

        left = ((self._slide_width - (img_size * len(images))) / 2) - 0.2

        for image in images:
            self._add_image_to_slide(
                image, slide, left, top, img_size
            )
            label_text = labeler(image)
            self._add_text_to_slide(
                slide, label_text, left, top + (img_size * 1.2), img_size, 1.5, 
                rotation=90)
            left += img_size
        return slide
    
    def _add_single_image_slide(self, image, title, metadata=None, img_scaler=0.65):
        '''General helper method for adding a slide with a single image to a
        presentation.

        :param image: Image to add to the slide
        :type image: Image
        :param title: Title to use for the slide
        :type title: str
        :param metadata: Additional information to write to the slide, defaults to None
        :type metadata: str, optional
        :param img_scaler: Scaler to apply to size of the image, defaults to 0.65
                            ,should be between 0 and 1. 1 is full sized image.
        :type img_scaler: float, optional
        :return: The new slide with Image added
        :rtype: slide
        '''
        new_slide = self._add_new_slide(5)
        new_slide.shapes.title.text = title
        img_size = self._slide_height * img_scaler
        self._add_image_to_slide(
            image, new_slide, self._bumper, 2, img_size)

        if metadata:
            metadata_offset = ((self._bumper + img_size) - self._slide_width) - self._bumper
            metadata_left = img_size + (self._bumper * 2)
            metadata_width = self._slide_width - metadata_left - self._bumper
            metadata_height = self._slide_height - 2 - self._bumper
            self._add_text_to_slide(
                new_slide, metadata, metadata_left, 2, metadata_width,
                metadata_height)
        return new_slide
        
    
    def _add_text_to_slide(self, slide, text, left, top, width, height,
                          rotation=0, font_size=14):
        '''Helper method to add text to a slide

        :param slide: Slide to add text to
        :type slide: slide
        :param text: Text to add to the slide
        :type text: str
        :param left: Left cordinate location of the text in inches
        :type left: float
        :param top: Top cordinate location of the text in inches
        :type top: float
        :param width: Width of the text in inches
        :type width: float
        :param height: Height of the text in inches
        :type height: float
        :param rotation: Rotation to apply to the text in degrees, defaults to 0
        :type rotation: int, optional
        :param font_size: Font size of text, defaults to 14
        :type font_size: int, optional
        :return: Slide with text added
        :rtype: slide
        '''
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height))
        text_box.rotation = rotation
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)

        return text_box
    
    def _add_image_to_slide(self, image, slide, left, top, height):
        '''Private method for adding images to a slide. If the :class:`Image`
        does not have a file written on the local machine as can be
        the case with saved runs who's image data only exists in
        their xtal files this method will write a temporary image
        file to the Polo :const:`TEMP_DIR` which then should be deleted after
        the presentaton file is written.

        :param image: Image to add to the slide
        :type image: Image
        :param slide: Slide to add the image to
        :type slide: slide
        :param left: Left cordinate location of the image in inches
        :type left: float
        :param top: Top cordinate location of the image in inches
        :type top: float
        :param height: Height of the image in inches
        :type height: float
        :return: []
        :rtype: [type]
        '''
        
        if os.path.isfile(image.path):
            img_path = image.path
        else:
            temp_path = str(TEMP_DIR.joinpath(str(hash(image.path))))
            with open(temp_path, 'wb') as tmp:
                tmp.write(base64.decodebytes(image.bites))
                self._temp_images.append(temp_path)
            img_path = temp_path
        
        return slide.shapes.add_picture(img_path, Inches(left), Inches(top), 
                                        height=Inches(height))
    
    # def add_classification_slide(self, well_number, rep_image):
    #     '''Add a slide containing details about an images MARCO
    #     and human classification in a table.

    #     :param well_number: Well number (index) of image to use in
    #                         the title of the slide 
    #     :type well_number: int
    #     :param rep_image: Image object to make slide from
    #     :type rep_image: Image
    #     '''
    #     new_slide = self._add_new_slide()
    #     title = 'Well {} Classifications'.format(well_number)
    #     new_slide.shapes.title.text = title

    #     data = [
    #         ['Human Classification', 'MARCO Classification']
    #         [rep_image.human_class, rep_image.machine_class]
    #     ]

    #     self.add_table_to_slide(new_slide, data, self._bumper, 2)
    #     # do for most recent human classification image if it exits


class BarTender():
    '''Class for organizing and accessing 
    :class:`~polo.utils.io_utils.Menu` data.

    :param cocktail_dir: Directory containing cocktail menu csv filepaths
    :type cocktail_dir: str or Path
    :param cocktail_meta: Path to cocktail metadata file which describes the contents of
                            each cocktail menu csv file
    :type cocktail_meta: Path or str

    Cocktail metadata file should be a csv file with the following
    headers ordered from top to bottom. Each header name is followed by a
    description.

    .. code-block:: text

        File Name: Name of cocktail menu file
        Dates Used: Range of dates the cocktail menu was used (m/d/y-m/d/y)
        Plate Number
        Screen Type: 'm' for membrane screens, 's' for soluble screens
    '''

    def __init__(self, cocktail_dir, cocktail_meta):
        self.cocktail_dir = cocktail_dir
        self.cocktail_meta = cocktail_meta
        self.menus = {}
        self.add_menus_from_metadata()

    @staticmethod
    def datetime_converter(date_string):
        '''General utility function for converting strings to datetime objects.
        Attempts to convert the string by trying a couple of datetime
        formats that are common in cocktail menu files and other
        locations in the HWI file universe Polo runs across.

        :param date_string: string to convert to datetime
        :type date_string: str
        :return: datetime object
        :rtype: datetime
        '''
        if date_string:
            date_string = str(date_string).strip()
            datetime_formats = ['%m/%d/%Y', '%m/%d/%y',
                                '%m-%d-%Y', '%m-%d-%y', '%Y-%m-%d %H:%M:%S']
            for form in datetime_formats:
                try:
                    return datetime.strptime(date_string, form)
                except ValueError as e:
                    continue

    @staticmethod
    def date_range_parser(date_range_string):
        '''Utility function for converting the date ranges in the cocktail
        metadata csv file to datetime objects using the 
        :meth:`~polo.utils.io_utils.BarTender.datetime_converter`
        method.

        Date ranges should have the format

        .. code-block:: text

            'start date - end date'

            If the date range is for the most recent cocktail menu then there
            will not be an end date and the format will be

            'start date - '

        :param date_range_string: string to pull dates out of   
        :type date_range_string: str
        :return: tuple of datetime objects, start date and end date
        :rtype: tuple
        '''
        s, e = date_range_string.split('-')
        s = BarTender.datetime_converter(s.strip())
        if e.strip():
            e = BarTender.datetime_converter(e.strip())
        else:
            e = None
        return s, e

    def add_menus_from_metadata(self):
        '''Adds :class:`Menu` objects to the :attr:`polo.utils.io_utils.BarTender.menus`
        attribute by reading the cocktail csv files included
        in the :const:`COCKTAIL_DATA_PATH` directory.'''
        if self.cocktail_meta:
            with open(str(self.cocktail_meta)) as menu_files:
                reader = csv.DictReader(menu_files)
                for row in reader:
                    path = os.path.join(
                        str(COCKTAIL_DATA_PATH), row['File Name'])
                    s, e = BarTender.date_range_parser(row['Dates Used'])
                    new_menu = Menu(path, s, e, row['Screen Type'])
                    self.menus[path] = new_menu
                    # add new menu to menus dict path to csv file is the menu
                    # key
        else:
            logger.warning('No cocktail metadata stored in {}'.format(self))

    def get_menu_by_date(self, date, type_='s'):
        '''Get a :class:`Menu` instance who's usage dates include the
        date passed through the `date` argument and
        matches the screen type passed through the `type_` argument.

        Screen types can either be 's' for 'soluble' screens or 'm' for
        membrane screens.

        :param date: Date to search menus with
        :type date: datetime
        :param type_: Type of screen to return (soluble or membrane)
        :type type_: str
        :return: menu matching the given date and type
        :rtype: Menu
        '''
        menus_keys_by_date = sorted(
                [key for key in self.menus.keys() if self.menus[key].type_ == type_],
                # keys matching only the specified type
                key=lambda key: self.menus[key].start_date
            )
        if isinstance(date, datetime):
            # search for a menu whos usage dates include this date
            # end up with a list of keys for menus of the specified type that
            # are sorted by the start date of their use at HWI cente
            for each_key in menus_keys_by_date:
                if date < self.menus[each_key].start_date:
                    return self.menus[each_key]
        return self.menus[menus_keys_by_date[-1]]

    def get_menus_by_type(self, type_='s'):
        '''Returns all :class:`Menu` instances of a given screen type.

        's' for soluble screens and 'm' for membrane screens. No other
        characters should be passed to `type_`.

        :param type_: Key for type of screen to return
        :type type_: str (max length 1)
        :return: list of menus of that screen type
        :rtype: list
        '''

        return [menu for menu in self.menus.values() if menu.type_ == type_]

    def get_menu_by_path(self, path):
        '''Returns a :class:`Menu` instance by its file path, which is
        used as the key for accessing the menus attribute normally.

        :param path: file path of a menu csv file
        :type path: str
        :return: Menu instance that is mapped to given path
        :rtype: Menu
        '''
        if path in self.menus:
            return self.menus[path]

    def get_menu_by_basename(self, basename):
        '''Uses the basename of a :class:`Menu` file path to return a :class:`~Menu` object.
        Useful for retrieving menus based on the text of comboBoxes since
        when menus are displayed to the user only the basename is used.

        :param basename: Basename of a :class:`~polo.utils.io_utils.Menu` file path
        :type basename: str
        :return: Menu instance who's basename matches the `basename` argument,
                returns None is no menu is found
        :rtype: Menu or None
        '''
        for menu_key in self.menus:  # self.menus is dictionary
            if os.path.basename(menu_key) == basename:
                return self.menus[menu_key]


class CocktailMenuReader():
    '''CocktailMenuReader instances should be used to read a csv file containing
    a collection of cocktail screens. The csv file should contain cocktail
    related formulations and assign each cocktail to a specific well in the
    screening plate.

    .. highlight:: python
    .. code-block:: python

        cocktail_menu = 'path/to/my/csv'
        my_reader = CocktailMenuReader(cocktail_menu)
        csv_cocktails = my_reader.read_menu_file()
        # csv_cocktails now holds list of Cocktail objects read from
        # cocktail_menu

    :param menu_file: Path to cocktail menu file to read. Should be csv
                        formated
    :type menu_file: str or Path 
    :param delim: Seperator for menu_file; really should not need to 
                    be changed, defaults to ','
    :type delim: str, optional
    '''

    cocktail_map = {  # map Cocktail attributes to index in menu rows
        0: 'well_assignment',
        1: 'number',
        8: 'pH',
        2: 'commercial_code'
    }
    formula_pos = 4  # each reagent could have a formula but only ever
    # one is included per cocktail entry
    # all other indicies are reagent names and concentrations

    def __init__(self, menu_file_path, delim=','):
        self.menu_file_path = menu_file_path
        self.delim = delim
        self._row_counter = 0

    @classmethod
    def set_cocktail_map(cls, map):
        '''Classmethod to edit the
        :attr:`~polo.utils.io_utils.CocktailMenuReader.cocktail_map`.
        The :attr:`~polo.utils.io_utils.CocktailMenuReader.cocktail_map`
        describes where the :class:`Cocktail` level information 
        is stored in a given cocktail row in the csv file. 
        It is a dictionary that maps specific indices in a row to
        :class:`Cocktail` attributes.

        The default cocktail_map dictionary is below.

        >>> cocktail_map = {
        0: 'well_assignment',
        1: 'number',
        8: 'pH',
        2: 'commercial_code'
        }

        This tells instances of CocktailMenuReader to look at index 0 of a row
        for the well_assignment attribute of the Cocktail class, index 1 for
        the number attribute of the Cocktail class, etc.

        :param map: Dictionary mapping csv row indices to Cocktail object
                    attributes
        :type map: dict
        '''
        cls.cocktail_map = map

    @classmethod
    def set_formula_pos(cls, pos):
        '''Classmethod to change the :attr:`CocktailMenuReader.formula_pos`
        attribute. The :attr:`~CocktailMenuReader.formula_pos`
        describes the location (base 0) of the chemical formula in a row of
        a cocktail menu file csv. For some reason, HWI cocktail menu files
        will only have one chemical formula per row (cocktail) no matter
        the number of reagents that composite that cocktail. This is why
        its location is represented using an int instead of a dict.

        Generally, :attr:`CocktailMenuReader.formula_pos`
        should not be changed without a very good
        reason as the position of the chemical formula is consistent across
        all HWI cocktail menu files.

        :param pos: Index where chemical formula can be found
        :type pos: int
        '''
        cls.formula_pos = pos

    def read_menu_file(self):
        '''Read the contents of cocktail menu csv file. The :class:`Menu` file path
        is read from the :attr:`Menu.menu_file_path` attribute. The first **two** lines
        of all the cocktail menu files included in Polo are header lines and
        so the reader will skip the first two lines before actually reading
        in any data. Each row is converted to a
        :class:`~polo.crystallography.cocktail.Cocktail` object and then
        added to a dictionary based on the  :class:`~polo.crystallography.cocktail.Cocktail`
        instance's well assignment.

        :return: Dictionary of Cocktail instances. Key value is the Cocktail's
                 well assignment in the screening plate (base 1).
        :rtype: dict
        '''
        cocktail_menu = {}
        with open(self.menu_file_path, 'r') as menu_file:
            reader = csv.reader(menu_file)
            next(reader)
            next(reader)  # skip first two rows
            for row in reader:
                d = {self.cocktail_map[index]: row[index]
                     for index in self.cocktail_map}
                new_cocktail = Cocktail(**d)
                new_cocktail.reagents = []
                # BUG: when new cocktail is initialized it contains the reagents
                # of all previously initialized cocktails. For now setting
                # new_cocktail.reagents to a new empty list fixes the
                # problem but does not address the source
                reagent_positions = [i for i in range(
                    len(row)) if i not in self.cocktail_map and i != self.formula_pos]


                for i in range(0, len(reagent_positions), 2):
                    try:
                        chem_add, con = (row[reagent_positions[i]],
                                        row[reagent_positions[i+1]])
                        if chem_add:
                            con = UnitValue.make_from_string(con)
                            new_cocktail.add_reagent(
                                Reagent(
                                    chemical_additive=chem_add,
                                    concentration=con
                                )
                            )
                    except IndexError as e:
                        # ETH 23
                        # Added this catch as a stop gap after adding in the new cocktail
                        # formulations. The csv file generated from the provided excel
                        # for this formulation has some differences between the old
                        # cocktail files that I have yet to determine that is causing
                        # there to be an odd number of reagent concentration paired
                        # columns. This function assumes for each chemical additive there
                        # is an associated concentration. So if a list with an odd number
                        # goes into this part of the function i+1 will throw an index error.
                        # Making this change allows the program to launch but I need to
                        # further investigate if the cocktails are being read correctly.
                        continue
                    
                cocktail_menu[int(new_cocktail.well_assignment)] = new_cocktail
        return cocktail_menu


class RunLinker():
    '''Class to hold methods relating to linking runs either by
    date or by spectrum.
    '''

    @staticmethod
    def link_visualizer(self, run):
        '''Create a text representation of the linked list structure between
        runs for debugging.

        :param run: Run to build representation from
        :type run: Run
        :return: Text representation of the linked list
        :rtype: str
        '''
        pass


    @staticmethod
    def the_big_link(runs):
        '''Wrapper method to do all the linking required for a collection of
        runs. First calls :meth:`~polo.utils.io_utils.RunLinker.unlink_runs_completely`
        to separate any existing links so things do not get tangled. Then 
        calls :meth:`~polo.utils.io_utils.RunLinker.link_runs_by_date` and
        :meth:`~polo.utils.io_utils.RunLinker.link_runs_by_spectrum`.

        :param runs: List of runs to link
        :type runs: list
        :return: List of runs with links made
        :rtype: list
        '''
        runs = RunLinker.unlink_runs_completely(runs)
        runs = RunLinker.link_runs_by_date(runs)
        runs = RunLinker.link_runs_by_spectrum(runs)

        return runs

    @staticmethod
    def link_runs_by_date(runs):
        # need to seperate out the runs that can be linked and not
        linkable_runs = [r for r in runs 
        if hasattr(r, 'link_to_next_date') and isinstance(r.date, datetime)]
        if linkable_runs:
            linkable_runs = [r for r in sorted(  # sort by date
            runs, key=lambda r: r.date) if r.image_spectrum == IMAGE_SPECS[0]]
            if linkable_runs:
                for i in range(0, len(linkable_runs)-1):
                    linkable_runs[i].link_to_next_date(linkable_runs[i+1])
                return list(set(runs).union(set(linkable_runs)))
        return runs

    @staticmethod
    def link_runs_by_spectrum(runs):
        '''Link a collection of :class:`~polo.crystallography.run.HWIRun` instances
        together by spectrum. All non-visible
        :class:`~polo.crystallography.run.HWIRun` instances
        are linked together in a monodirectional circular linked list.
        Each visible :class:`~polo.crystallography.run.HWIRun` instance
        will then point to the same non-visible run through
        their :attr:`~polo.crystallography.run.HWIRun.alt_spectrum`
        attribute as a way to access the non-visible
        linked list.

        :param runs: List of runs to link together
        :type runs: list
        :return: List of runs linked by spectrum
        :rtype: list
        '''
        # for now this links all runs of the sample to the alt spectrums when
        linkable_runs = [r for r in runs if hasattr(r, 'link_to_alt_spectrum')]
        visible, other = [], []
        for run in linkable_runs:
            if run.image_spectrum == IMAGE_SPECS[0]:  # visible images only
                visible.append(run)
            else:
                other.append(run)
        if other:
            if len(other) > 1:
                other = sorted(other, key=lambda o: len(str(o.image_spectrum)))
                for i in range(len(other)-1):
                    other[i].link_to_alt_spectrum(other[i+1])
                other[-1].link_to_alt_spectrum(other[0])

            if visible:
                for run in visible:
                    run.link_to_alt_spectrum(other[0])  # link to first run of alts
        return runs

    @staticmethod
    def unlink_runs_completely(runs):
        '''Cuts all links between the :class:`~polo.crystallography.run.HWIRun` instances
        passed through the `runs` argument and the
        :class:`~polo.crystallography.image.Image` instances
        in those runs.

        :param runs: List of runs
        :type runs: list
        :return: List of runs without any links
        :rtype: list
        '''
        for i, _ in enumerate(runs):
            runs[i].previous_run, runs[i].next_run, runs[i].alt_spectrum = None, None, None
            for image in runs[i].images:
                image.next_image, image.previous_image, image.alt_image = (
                    None, None, None
                )
        return runs


class XmlReader():

    #platedef_key = 'platedef'  # keyword that is always in plate definition
    # xml file names

    def __init__(self, xml_path=None):
        '''XmlReader class can be used to read the xml metadata files that are
        included in HWI screening run rar archives. Currently, is primarily meant
        to extract metadata about the plate and the sample in that plate.

        :param xml_path: File path to xml file
        :type xml_path: str or Path
        :param xml_files: list of xml file paths, defaults to []
        :type xml_files: list, optional
        '''
        self.xml_path = xml_path
        self._tree = ET.parse(str(xml_path))
        self._root = self._tree.getroot()
        logger.debug('Made {} for {}'.format(self, str(self.xml_path)))


    @staticmethod
    def get_data_from_xml_element(xml_element):
        '''Return the data stored in an `xml_element`. Helper method
        for reading xml files.

        :param xml_element: xml element to read data from
        :type xml_element: [type]
        :return: Dictionary of data stored in xml element
        :rtype: dict
        '''
        return {elem.tag: elem.text for elem in xml_element
                if elem.tag and elem.text}
    
    def __getitem__(self, index):
        return self.get_data_from_xml_element(self._root[index])


class Menu():  # holds the dictionary of cocktails

    def __init__(self, path, start_date, end_date, type_, cocktails={}):
        '''Creates a :class:`Menu` instance. :class:`Menu` objects are used to organize a single
        screening run plate set up. They should contain 1536 unique screening
        conditions; one for each well in the HWI high-throughput plate. HWI
        has altered what cocktails are used in each of the 1536 wells over the
        years and so many versions of cocktail menus have been used. A single
        :class:`Menu` instance represents one of these versions and accordingly has a
        start and end date to identify when the menu was used. Additionally, 
        HWI offers two types of high throughput screens; membrane or soluble
        screens. Both use 1536 well plates but have very different chemical
        conditions at the same well index. 

        :class:`Menu`s that hold conditions for soluble screens are the `type_`
        attribute set to 's' and menus that hold conditions for membrane
        screens have the `type_` attribute set to 'm'.

        :param path: Location of the menu csv file
        :type path: str
        :param start_date: Date when this screen menu was first used
        :type start_date: datetime
        :param end_date: Last date this screen menu was used
        :type end_date: datetime
        :param type_: Membrane or soluble screen
        :type type_: str
        '''
        self.start_date = start_date
        self.end_date = end_date
        self.type_ = type_
        self.path = path
        self.cocktails = cocktails  # holds all cocktails (items on the menu)
        logger.debug('Created Menu for path {}'.format(self.path))

    @property
    def cocktails(self):
        return self._cocktails

    @cocktails.setter
    def cocktails(self, new_cocktails):
        if new_cocktails:
            self._cocktails = new_cocktails
        else:
            self._cocktails = CocktailMenuReader(self.path).read_menu_file()
            logger.debug('Added {} Cocktails'.format(len(self._cocktails)))

    def __len__(self):
        if self.cocktails:
            return len(self.cocktails)
        else:
            return None
    
    def __getitem__(self, key):
        return self.cocktails[key]
    

# orphan functions that have not made it into a class yet
# =============================================================================


def write_screen_html(plate_list, well_number, run_name, x_reagent,
                      y_reagent, well_volume, output_path):

    template = SCREEN_HTML_TEMPLATE
    template_string = open(str(SCREEN_HTML_TEMPLATE)).read()
    template = Template(template_string)

    html = template.render(plate_list=plate_list, run_name=run_name,
                           well_number=well_number, date=datetime.now(),
                           x_reagent_stock=x_reagent,
                           y_reagent_stock=y_reagent,
                           well_volume=str(well_volume))

    output_path = Path(str(output_path))
    if output_path.suffix != '.html':
        output_path = str(output_path.with_suffix('.html'))

    with open(str(output_path), 'w') as screen_html:
        screen_html.write(html)


def parse_HWI_filename_meta(HWI_image_file):
    '''HWI images have a standard file nameing schema that gives info about when
    they are taken and well number and that kind of thing. This function returns
    that data
    '''
    HWI_image_file = os.path.basename(HWI_image_file)
    return (
        HWI_image_file[:10],
        int(HWI_image_file[10:14].lstrip('0')),
        datetime.strptime(HWI_image_file[14:22], '%Y''%m''%d'),
        HWI_image_file[22:]
    )


def list_dir_abs(parent_dir, allowed=False):
    files, allowed_files = os.listdir(parent_dir), []
    for f in files:
        abs_file_path = os.path.join(parent_dir, f)
        if allowed:
            ext = os.path.splitext(f)[-1]
            if ext in ALLOWED_IMAGE_TYPES:
                allowed_files.append(Path(abs_file_path))
        else:
            allowed_files.append(Path(abs_file_path))

    return allowed_files


def parse_hwi_dir_metadata(dir_name):
    try:
        dir_name = os.path.basename(dir_name)
        image_type = dir_name.split('-')[-1]
        plate_id = dir_name[:10]
        date = datetime.strptime(dir_name[10:].split('-')[0],
                                 '%Y''%m''%d''%H''%M')

        return image_type, plate_id, date, dir_name  # last is suggeseted run name
    except ValueError as e:
        logger.error('Caught {} at {} attempting to parse {}'.format(
            e, parse_hwi_dir_metadata, dir_name
        ))
        return e


def directory_validator(dir_path):
    if os.path.exists(dir_path):
        if os.path.isdir(dir_path):
            files = list_dir_abs(dir_path, allowed=True)
            if files:
                return True
            else:
                e = ForbiddenImageTypeError
                logger.warning('Directory validation failed with {}'.format(e))
                return e
        else:
            e = NotADirectoryError
            logger.warning('Directory validation failed with {}'.format(e))
            return e
    else:
        e = FileNotFoundError
        logger.warning('Directory validation failed with {}'.format(e))
        return FileNotFoundError


def check_for_missing_images(dir_path, expected_image_count):
    number_images = len(list_dir_abs(dir_path, allowed=True))
    if number_images != expected_image_count:
        return False
    else:
        return True


def run_name_validator(new_run_name, current_run_names):
    # checks that new run name is utf 8 and
    try:
        new_run_name = new_run_name.encode('utf-8')
        new_run_name = new_run_name.decode('utf-8')
    except UnicodeError as e:
        logger.warning('{} failed run name validation with {}'.format(
            new_run_name, e
        ))
        return e
    if new_run_name == '' or new_run_name == None:
        return TypeError
    if new_run_name not in current_run_names:
        return True
    else:
        return False


def if_dir_not_exists_make(parent_dir, child_dir=None):
    '''If only parent_dir is given attempts to make that directory. If parent
    and child are given tries to make a directory child_dir within parent dir.
    '''
    if child_dir:
        path = os.path.join(parent_dir, child_dir)
    else:
        path = parent_dir

    if not os.path.exists(path):
        try:
            os.mkdir(path)
        except FileNotFoundError as e:
            logger.error('Could not make directory with path {}. Returned {}'.format(
                path, e
            ))
            return e

    return path

# NOTE imports moved to end as quick fix. Files imported here import classes
# and functions from io_utils so entire io_utils script must be run before
# other scripts are run to avoid import errors
# RUN_TYPES is also here because until those scripts are run they are not
# included in sys.modules dictionary

from polo.crystallography.cocktail import *
from polo.crystallography.image import Image
from polo.crystallography.run import *
RUN_TYPES = sorted(
        [types[-1] for types in 
        inspect.getmembers(sys.modules['polo.crystallography.run'], inspect.isclass)
        if issubclass(types[-1], Run)],
        key=lambda c: c.import_priority,
        reverse=True)